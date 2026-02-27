"""
Generate a professional 2-slide PPTX deck for Compliance Copilot.
Slide 1: Business Value Proposition
Slide 2: Architecture + Azure Integration
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Constants
DARK_BG = RGBColor(0x0D, 0x11, 0x17)
CARD_BG = RGBColor(0x16, 0x1B, 0x22)
ACCENT_BLUE = RGBColor(0x58, 0xA6, 0xFF)
ACCENT_GREEN = RGBColor(0x3F, 0xB9, 0x50)
ACCENT_YELLOW = RGBColor(0xD2, 0x99, 0x22)
ACCENT_RED = RGBColor(0xF8, 0x51, 0x49)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x8B, 0x94, 0x9E)
LIGHT_GRAY = RGBColor(0xC9, 0xD1, 0xD9)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Use blank layout
blank_layout = prs.slide_layouts[6]


def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, border_color=None, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    if radius:
        shape.adjustments[0] = radius
    return shape


def add_text(slide, left, top, width, height, text, font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_metric_card(slide, left, top, before_val, after_val, label, accent_color):
    card_w = Inches(3.6)
    card_h = Inches(1.8)
    add_rect(slide, left, top, card_w, card_h, CARD_BG, RGBColor(0x30, 0x36, 0x3D), 0.04)

    # Label
    add_text(slide, left + Inches(0.3), top + Inches(0.2), card_w - Inches(0.6), Inches(0.4),
             label, font_size=11, color=GRAY, bold=False)

    # Before → After
    txBox = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.6), card_w - Inches(0.6), Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = before_val
    run1.font.size = Pt(14)
    run1.font.color.rgb = GRAY
    run1.font.name = "Segoe UI"
    run1.font.strikethrough = True

    run2 = p.add_run()
    run2.text = "  →  "
    run2.font.size = Pt(14)
    run2.font.color.rgb = GRAY
    run2.font.name = "Segoe UI"

    run3 = p.add_run()
    run3.text = after_val
    run3.font.size = Pt(28)
    run3.font.color.rgb = accent_color
    run3.font.bold = True
    run3.font.name = "Segoe UI"


# ═══════════════════════════════════════════
# SLIDE 1: Business Value Proposition
# ═══════════════════════════════════════════
slide1 = prs.slides.add_slide(blank_layout)
add_bg(slide1)

# Top accent line
accent_line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.06))
accent_line.fill.solid()
accent_line.fill.fore_color.rgb = ACCENT_BLUE
accent_line.line.fill.background()

# Title area
add_text(slide1, Inches(0.8), Inches(0.4), Inches(8), Inches(0.4),
         "COMPLIANCE COPILOT", font_size=12, color=ACCENT_BLUE, bold=True)

add_text(slide1, Inches(0.8), Inches(0.75), Inches(10), Inches(0.8),
         "AI-Powered Compliance Review for Every Pull Request",
         font_size=32, color=WHITE, bold=True)

add_text(slide1, Inches(0.8), Inches(1.5), Inches(10), Inches(0.6),
         "Built with GitHub Copilot SDK  |  Work IQ  |  Fabric IQ  |  Azure",
         font_size=14, color=GRAY)

# Horizontal divider
div = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.2), Inches(11.7), Pt(1))
div.fill.solid()
div.fill.fore_color.rgb = RGBColor(0x30, 0x36, 0x3D)
div.line.fill.background()

# Metric cards row
add_metric_card(slide1, Inches(0.8), Inches(2.5), "3-5 days", "30 sec", "Review Time", ACCENT_GREEN)
add_metric_card(slide1, Inches(4.8), Inches(2.5), "~60%", "95%", "Violation Detection", ACCENT_BLUE)
add_metric_card(slide1, Inches(8.8), Inches(2.5), "Manual", "100%", "Audit Coverage", ACCENT_YELLOW)

# Problem → Solution section
add_text(slide1, Inches(0.8), Inches(4.6), Inches(5.5), Inches(0.3),
         "THE PROBLEM", font_size=11, color=ACCENT_RED, bold=True)

problem_text = (
    "Manual compliance reviews create a 3-5 day bottleneck on every pull request. "
    "Engineering teams either skip reviews (risk) or wait days for approval (velocity loss). "
    "Audit trails are incomplete, inconsistent, and not exportable."
)
add_text(slide1, Inches(0.8), Inches(4.95), Inches(5.5), Inches(1.2),
         problem_text, font_size=13, color=LIGHT_GRAY)

add_text(slide1, Inches(7.0), Inches(4.6), Inches(5.5), Inches(0.3),
         "THE SOLUTION", font_size=11, color=ACCENT_GREEN, bold=True)

solution_text = (
    "Compliance Copilot uses the GitHub Copilot SDK to automatically review every PR "
    "against SOC 2 Type II and HIPAA controls. The agent posts structured findings with "
    "severity ratings, control references, and remediation guidance. Every review is "
    "recorded in Fabric IQ for audit dashboards and compliance trending."
)
add_text(slide1, Inches(7.0), Inches(4.95), Inches(5.5), Inches(1.4),
         solution_text, font_size=13, color=LIGHT_GRAY)

# Bottom: How it works (4 steps)
add_text(slide1, Inches(0.8), Inches(6.4), Inches(10), Inches(0.3),
         "HOW IT WORKS", font_size=11, color=ACCENT_BLUE, bold=True)

steps = [
    ("1", "PR Opened", "Webhook triggers agent"),
    ("2", "Agent Analyzes", "5 custom tools via Copilot SDK"),
    ("3", "Review Posted", "Structured findings + score"),
    ("4", "Audit Stored", "Fabric IQ dashboard + export"),
]
for i, (num, title, desc) in enumerate(steps):
    x = Inches(0.8 + i * 3.1)
    # Step number circle
    circle = slide1.shapes.add_shape(MSO_SHAPE.OVAL, x, Inches(6.75), Inches(0.35), Inches(0.35))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT_BLUE
    circle.line.fill.background()
    circle.text_frame.paragraphs[0].text = num
    circle.text_frame.paragraphs[0].font.size = Pt(12)
    circle.text_frame.paragraphs[0].font.color.rgb = WHITE
    circle.text_frame.paragraphs[0].font.bold = True
    circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    circle.text_frame.paragraphs[0].font.name = "Segoe UI"

    add_text(slide1, x + Inches(0.5), Inches(6.72), Inches(2.4), Inches(0.3),
             title, font_size=13, color=WHITE, bold=True)
    add_text(slide1, x + Inches(0.5), Inches(6.98), Inches(2.4), Inches(0.3),
             desc, font_size=11, color=GRAY)


# ═══════════════════════════════════════════
# SLIDE 2: Architecture + Azure Integration
# ═══════════════════════════════════════════
slide2 = prs.slides.add_slide(blank_layout)
add_bg(slide2)

# Top accent line
accent_line2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.06))
accent_line2.fill.solid()
accent_line2.fill.fore_color.rgb = ACCENT_BLUE
accent_line2.line.fill.background()

add_text(slide2, Inches(0.8), Inches(0.4), Inches(8), Inches(0.4),
         "ARCHITECTURE", font_size=12, color=ACCENT_BLUE, bold=True)

add_text(slide2, Inches(0.8), Inches(0.75), Inches(10), Inches(0.7),
         "Azure-First Enterprise Architecture",
         font_size=28, color=WHITE, bold=True)

# Architecture diagram area (left side, main)
arch_card = add_rect(slide2, Inches(0.8), Inches(1.6), Inches(7.8), Inches(5.2), CARD_BG, RGBColor(0x30, 0x36, 0x3D), 0.02)

# Architecture flow (text-based diagram inside the card)
arch_lines = [
    ("GitHub PR", ACCENT_BLUE, 2.0),
    ("     │  webhook (HMAC-SHA256)", GRAY, 2.35),
    ("     ▼", GRAY, 2.6),
    ("Express Server  (Azure Container Apps)", ACCENT_GREEN, 2.85),
    ("     │", GRAY, 3.1),
    ("     ▼", GRAY, 3.3),
    ("Copilot SDK Agent  (gpt-4.1)", ACCENT_YELLOW, 3.55),
    ("  ┌──────┬──────┬──────┬──────┐", GRAY, 3.85),
    ("  │      │      │      │      │", GRAY, 4.05),
]

# Use monospace font for the diagram
for text, color, y_pos in arch_lines:
    add_text(slide2, Inches(1.8), Inches(y_pos), Inches(6), Inches(0.25),
             text, font_size=12, color=color, font_name="Consolas")

# Tool boxes row
tools = [
    ("fetch_pr_diff", "GitHub API"),
    ("query_policies", "Work IQ"),
    ("classify_data", "Azure Purview"),
    ("check_exceptions", "Work IQ"),
    ("store_audit", "Fabric IQ"),
]
for i, (tool, service) in enumerate(tools):
    x = Inches(1.2 + i * 1.45)
    tool_card = add_rect(slide2, x, Inches(4.3), Inches(1.35), Inches(0.95), RGBColor(0x21, 0x26, 0x2D), ACCENT_BLUE, 0.06)
    add_text(slide2, x + Inches(0.08), Inches(4.35), Inches(1.2), Inches(0.45),
             tool, font_size=8, color=WHITE, bold=True, font_name="Consolas")
    add_text(slide2, x + Inches(0.08), Inches(4.75), Inches(1.2), Inches(0.3),
             service, font_size=8, color=ACCENT_BLUE)

# Output row
add_text(slide2, Inches(1.8), Inches(5.4), Inches(6), Inches(0.25),
         "  │      │      │      │      │", font_size=12, color=GRAY, font_name="Consolas")
add_text(slide2, Inches(1.8), Inches(5.6), Inches(6), Inches(0.25),
         "  └──────┴──────┴──────┴──────┘", font_size=12, color=GRAY, font_name="Consolas")
add_text(slide2, Inches(1.8), Inches(5.85), Inches(6), Inches(0.25),
         "     ▼                    ▼", font_size=12, color=GRAY, font_name="Consolas")

# Output boxes
out1 = add_rect(slide2, Inches(1.5), Inches(6.15), Inches(2.8), Inches(0.55), RGBColor(0x21, 0x26, 0x2D), ACCENT_GREEN, 0.06)
add_text(slide2, Inches(1.6), Inches(6.2), Inches(2.6), Inches(0.45),
         "PR Comment (score + findings)", font_size=10, color=ACCENT_GREEN, bold=True)

out2 = add_rect(slide2, Inches(4.8), Inches(6.15), Inches(2.8), Inches(0.55), RGBColor(0x21, 0x26, 0x2D), ACCENT_YELLOW, 0.06)
add_text(slide2, Inches(4.9), Inches(6.2), Inches(2.6), Inches(0.45),
         "Fabric IQ Dashboard + Audit", font_size=10, color=ACCENT_YELLOW, bold=True)

# Right side: Azure integration stack
add_text(slide2, Inches(9.0), Inches(1.6), Inches(4), Inches(0.3),
         "AZURE INTEGRATION", font_size=11, color=ACCENT_BLUE, bold=True)

azure_services = [
    ("Azure Container Apps", "Deployment & scaling", ACCENT_GREEN),
    ("Entra ID (Azure AD)", "Authentication & RBAC", ACCENT_BLUE),
    ("Azure Key Vault", "Secrets management", ACCENT_YELLOW),
    ("Azure Purview", "Data classification (PII/PHI)", RGBColor(0xA7, 0x8B, 0xFA)),
    ("Azure Monitor", "Telemetry & alerting", ACCENT_RED),
    ("Work IQ", "Policy management", ACCENT_BLUE),
    ("Fabric IQ", "Audit trail & dashboards", ACCENT_GREEN),
]

for i, (name, desc, color) in enumerate(azure_services):
    y = Inches(2.05 + i * 0.62)
    # Color indicator bar
    bar = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.0), y, Inches(0.08), Inches(0.45))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    add_text(slide2, Inches(9.2), y - Inches(0.02), Inches(3.5), Inches(0.28),
             name, font_size=12, color=WHITE, bold=True)
    add_text(slide2, Inches(9.2), y + Inches(0.22), Inches(3.5), Inches(0.25),
             desc, font_size=10, color=GRAY)

# Compliance frameworks
add_text(slide2, Inches(9.0), Inches(6.4), Inches(4), Inches(0.3),
         "FRAMEWORKS", font_size=11, color=ACCENT_BLUE, bold=True)

fw_text = "SOC 2 Type II  ·  HIPAA  ·  Extensible"
add_text(slide2, Inches(9.0), Inches(6.7), Inches(4), Inches(0.3),
         fw_text, font_size=12, color=LIGHT_GRAY)

# Repo link at bottom
add_text(slide2, Inches(0.8), Inches(7.05), Inches(12), Inches(0.35),
         "github.com/adgranoff/Compliance-Copilot",
         font_size=11, color=ACCENT_BLUE, alignment=PP_ALIGN.RIGHT)


# Save
output_path = r"C:\projects\contest\presentations\ComplianceCopilot.pptx"
prs.save(output_path)
print(f"PPTX saved to {output_path}")
